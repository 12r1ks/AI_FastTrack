"""Run: uv run --with python-pptx python scrn_Presentation/make_pptx.py"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x11, 0x17)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xC9, 0xD1, 0xD9)
BLUE    = RGBColor(0x58, 0xA6, 0xFF)
GREEN   = RGBColor(0x3F, 0xB9, 0x50)
SUBTLE  = RGBColor(0x6E, 0x76, 0x81)
TH_BG   = RGBColor(0x21, 0x26, 0x2D)
TD_BG   = RGBColor(0x16, 0x1B, 0x22)

HERE = Path(__file__).parent

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = BG
    return s

def tb(s, text, l, t, w, h, size=14, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def title_block(s, text, sub=None):
    tb(s, text, Inches(0.7), Inches(0.2), Inches(11.9), Inches(0.75),
       size=30, bold=True, color=WHITE)
    if sub:
        tb(s, sub, Inches(0.7), Inches(0.88), Inches(11.9), Inches(0.38),
           size=14, color=BLUE)

def block(s, head, items, l, t, w, h, size=12,
          head_color=BLUE, item_color=GRAY):
    box = s.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    entries = [('h', head)] + [('i', x) for x in items]
    for kind, text in entries:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size  = Pt(size + 1 if kind == 'h' else size)
        r.font.bold  = (kind == 'h')
        r.font.color.rgb = head_color if kind == 'h' else item_color

def table(s, headers, rows, l, t, w, h, col_widths=None):
    shape = s.shapes.add_table(len(rows) + 1, len(headers), l, t, w, h)
    tbl   = shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TH_BG
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = hdr
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE

    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TD_BG
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = val
            r.font.size = Pt(11)
            r.font.color.rgb = GRAY

def pic(s, name, l, t, w=None, h=None):
    path = str(HERE / name)
    kw = {}
    if w: kw['width']  = w
    if h: kw['height'] = h
    s.shapes.add_picture(path, l, t, **kw)

# ── Slide 1 — Title ───────────────────────────────────────────────────────────
s1 = new_slide()
tb(s1, "CityPark Assistant",
   Inches(0.7), Inches(2.0), Inches(11.9), Inches(1.2),
   size=44, bold=True, align=PP_ALIGN.CENTER)
tb(s1, "Stage 1 — RAG Chatbot",
   Inches(0.7), Inches(3.1), Inches(11.9), Inches(0.6),
   size=22, color=BLUE, align=PP_ALIGN.CENTER)
tb(s1, "[ Name  ·  Date  ·  Context ]",
   Inches(0.7), Inches(4.0), Inches(11.9), Inches(0.4),
   size=14, italic=True, color=SUBTLE, align=PP_ALIGN.CENTER)

# ── Slide 2 — Architecture ────────────────────────────────────────────────────
s2 = new_slide()
title_block(s2, "Architecture", "LangGraph StateGraph")
pic(s2, "LangSmith.png", Inches(0.5), Inches(1.4), h=Inches(5.3))
block(s2, "Nodes",
    [
        "init  —  resets loop counter per request",
        "guard_route  —  LLM classifier → llm_call / reject",
        "llm_call  —  main node, tool-bound",
        "tool_node  —  dispatches tools, updates state",
        "pre_end_guard  —  output filter before response leaves",
        "reject  —  canned response, no reasoning exposed",
    ],
    Inches(6.5), Inches(1.5), Inches(6.4), Inches(5.2),
    size=12, head_color=BLUE)

# ── Slide 3 — Guardrails ──────────────────────────────────────────────────────
s3 = new_slide()
title_block(s3, "Guardrails", "Three independent layers")

block(s3, "Input Guard  (guard_route)",
    [
        "Structured-output LLM classifier on every message",
        "Routes to reject: off-topic, prompt injection, adversarial roleplay",
        "Canned response — no internal reasoning exposed",
    ],
    Inches(0.7), Inches(1.5), Inches(5.8), Inches(2.2), size=12)

block(s3, "Output Guard  (pre_end_guard)",
    [
        "Regex pass on every response before it leaves the graph",
        "Masks phone numbers:  +37529…7600  →  +37*****7600",
        "Runs on both normal and reject-path responses",
    ],
    Inches(0.7), Inches(3.6), Inches(5.8), Inches(2.2), size=12)

block(s3, "Data Isolation",
    [
        "FAISS  —  static parking info only  (no PII)",
        "SQLite  —  spot inventory + booking slots  (no names, no phones)",
        "Personal data lives only in graph state  (in-memory, session-scoped)",
        "Tools have no path to personal data by design",
        "Input and output guards cover different threat vectors:",
        "  semantic (LLM-based)  vs  structural (regex-based)",
    ],
    Inches(7.0), Inches(1.5), Inches(5.9), Inches(4.5),
    size=12, head_color=GREEN)

# ── Slide 4 — Technical Decisions ────────────────────────────────────────────
s4 = new_slide()
title_block(s4, "Key Technical Decisions")

decisions = [
    ("LangGraph StateGraph",
     "Explicit state, conditional routing, and interrupt/resume support for Stage 2+"),
    ("LlamaIndex FAISS",
     "Direct embedding config control; avoids langchain-community abstraction layer"),
    ("InMemorySaver",
     "Sufficient for this project — fully supports interrupt/resume within the same process"),
    ("Dual compile pattern",
     "agent.py compiles without checkpointer for eval; app.py recompiles with InMemorySaver for HTTP"),
    ("[ Your decision ]",
     "[ Why you made it ]"),
]

for i, (dec, reason) in enumerate(decisions):
    top = Inches(1.45) + i * Inches(1.08)
    tb(s4, dec,    Inches(0.7), top,           Inches(4.8), Inches(0.45), size=13, bold=True,  color=BLUE)
    tb(s4, reason, Inches(5.7), top + Inches(0.02), Inches(7.2), Inches(0.55), size=12, color=GRAY)

# ── Slide 5 — Tools ───────────────────────────────────────────────────────────
s5 = new_slide()
title_block(s5, "Tools & Data Sources")
table(s5,
    ["Tool", "When the LLM calls it", "Data source"],
    [
        ["Price_Calculator",
         "User needs a cost estimate",
         "Pure computation"],
        ["Retrieve_data_from_company_database",
         "Questions about locations, hours, policy, rates",
         "FAISS vector store"],
        ["query_available_spots",
         "Checking if a time slot is free",
         "SQLite — spots + bookings"],
        ["store_or_update_info_for_parking_proposal",
         "All reservation fields have been collected",
         "Graph state (MessagesState)"],
    ],
    Inches(0.7), Inches(1.4), Inches(11.9), Inches(4.8),
    col_widths=[3.5, 5.0, 3.0])

# ── Slide 6 — Eval Results ────────────────────────────────────────────────────
s6 = new_slide()
title_block(s6, "RAG Evaluation", "RAGAS — 10 questions · 5 metrics")
pic(s6, "Evals report.png", Inches(0.4), Inches(1.38), w=Inches(9.3))
block(s6, "Metric guide",
    [
        "Precision  0.93  — right chunks retrieved",
        "Recall  1.00  — no relevant info missed",
        "Faithfulness  1.00  — no hallucinations",
        "Relevancy  0.94  — responses on-topic",
        "Noise Sensitivity  0.44  — lower is better",
        "All 10: polite",
        "",
        "[ Your observations ]",
    ],
    Inches(9.9), Inches(1.45), Inches(3.05), Inches(5.6),
    size=11, head_color=BLUE)

# ── Slide 7 — Demo ────────────────────────────────────────────────────────────
s7 = new_slide()
title_block(s7, "Demo", "CityPark Assistant — web UI")
pic(s7, "Web-app.png", Inches(2.5), Inches(1.25), h=Inches(5.5))
tb(s7, "[ Walk through a reservation flow here ]",
   Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
   size=12, italic=True, color=SUBTLE, align=PP_ALIGN.CENTER)

# ── Slide 8 — What's Next ─────────────────────────────────────────────────────
s8 = new_slide()
title_block(s8, "What's Next", "Stages 2 – 4")

stages = [
    ("Stage 2  —  Human-in-the-Loop",
     ["Admin approval via LangGraph interrupt()",
      "Graph pauses, admin reviews proposal",
      "Resumes on accept / reject decision",
      "Admin UI endpoint"]),
    ("Stage 3  —  MCP Server",
     ["Approved reservations written via MCP",
      "Confirmed bookings persisted",
      "[ Your description ]"]),
    ("Stage 4  —  Orchestration",
     ["Multi-agent coordination",
      "[ Your description ]"]),
]

for i, (stitle, items) in enumerate(stages):
    l = Inches(0.7) + i * Inches(4.2)
    tb(s8, stitle, l, Inches(1.4), Inches(3.9), Inches(0.5),
       size=13, bold=True, color=BLUE)
    box = s8.shapes.add_textbox(l, Inches(1.95), Inches(3.9), Inches(4.5))
    tf  = box.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = item
        r.font.size = Pt(12)
        r.font.color.rgb = SUBTLE if item.startswith("[") else GRAY

# ── Save ──────────────────────────────────────────────────────────────────────
out = HERE / "CityPark_Stage1.pptx"
prs.save(str(out))
print(f"Saved: {out}")
